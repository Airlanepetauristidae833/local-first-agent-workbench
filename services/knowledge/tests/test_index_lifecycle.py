from __future__ import annotations

import asyncio
import json
import os
import sqlite3
import time
import zipfile
from pathlib import Path
from uuid import uuid4

import numpy as np
import pytest

from app import (
    CHUNK_VERSION,
    DOCUMENT_EXTENSIONS,
    CodexHandoff,
    CodexHandoffResult,
    KnowledgeService,
    ProgressCapture,
    ProjectCreate,
    ProjectDelete,
    Query,
    ResearchCapture,
    ResearchSource,
)


class FakeModel:
    def __init__(self) -> None:
        self.fail = False

    def encode(self, texts, **_kwargs):
        if self.fail:
            raise RuntimeError("embedding failed")
        vectors = []
        for text in texts:
            value = float((sum(text.encode("utf-8")) % 97) + 1)
            vectors.append([value, value / 2, value / 3])
        vectors = np.asarray(vectors, dtype=np.float32)
        norms = np.linalg.norm(vectors, axis=1, keepdims=True)
        return vectors / norms


def make_service(tmp_path: Path, *, max_text_chars: int = 500_000) -> KnowledgeService:
    return KnowledgeService(
        data=tmp_path / "data",
        vault=tmp_path / "vault",
        sources=tmp_path / "sources",
        model=FakeModel(),
        embedding_model="test/embedding",
        embedding_revision="revision-1",
        max_text_chars=max_text_chars,
    )


def create_project(service: KnowledgeService, identifier: str = "sample") -> Path:
    service.add(ProjectCreate(id=identifier, name="Sample", source_paths=[]))
    return service.safe(service.vault, f"Projects/AI Workbench/{identifier}")


def test_manifest_detects_add_change_delete_and_records_file_stats(
    tmp_path: Path,
) -> None:
    service = make_service(tmp_path, max_text_chars=80)
    vault = create_project(service)
    (vault / "Sources").mkdir()
    (vault / "Sources" / "main.swift").write_text(
        'struct App { let title = "first" }', encoding="utf-8"
    )
    (vault / "Sources" / "web.ts").write_text(
        "export const answer: number = 42", encoding="utf-8"
    )
    (vault / "Sources" / "project.pbxproj").write_text(
        "PRODUCT_NAME = Sample;", encoding="utf-8"
    )
    (vault / "long.txt").write_text("x" * 200, encoding="utf-8")
    (vault / "empty.md").write_text("", encoding="utf-8")
    (vault / "asset.bin").write_bytes(b"binary")
    (vault / "node_modules").mkdir()
    (vault / "node_modules" / "ignored.js").write_text(
        "should not be indexed", encoding="utf-8"
    )

    first = asyncio.run(service.index("sample", reason="test"))
    project = service.project("sample")
    manifest = project["index_manifest"]
    stats = project["index_stats"]
    sources = {Path(item["source"]).name for item in manifest["files"]}

    assert first["rebuilt"] is True
    assert manifest["chunk_version"] == CHUNK_VERSION
    assert manifest["embedding_revision"] == "revision-1"
    assert {"main.swift", "web.ts", "project.pbxproj", "long.txt"} <= sources
    assert "ignored.js" not in sources
    assert stats["unsupported_files"] == 1
    assert stats["truncated_files"][0]["source"].endswith("long.txt")
    assert stats["failed_files"] == 0
    assert "files" not in service.public_project(project)["index_manifest"]
    assert service.public_project(project)["index_stats"]["indexed_files"] >= 1

    current = asyncio.run(service.index_if_stale("sample", reason="test sweep"))
    assert current["rebuilt"] is False

    old_fingerprint = manifest["fingerprint"]
    (vault / "Sources" / "web.ts").unlink()
    rebuilt = asyncio.run(service.index_if_stale("sample", reason="delete"))
    assert rebuilt["rebuilt"] is True
    assert service.project("sample")["index_manifest"]["fingerprint"] != old_fingerprint
    assert all(
        not item["source"].endswith("web.ts")
        for item in service.project("sample")["index_manifest"]["files"]
    )


def test_global_search_returns_ranked_provenance_across_projects(tmp_path: Path) -> None:
    service = make_service(tmp_path)
    first = create_project(service, "first")
    second = create_project(service, "second")
    (first / "evidence.md").write_text("alpha evidence", encoding="utf-8")
    (second / "evidence.md").write_text("beta evidence", encoding="utf-8")
    asyncio.run(service.index("first", reason="test"))
    asyncio.run(service.index("second", reason="test"))

    result = service.search_all(Query(query="evidence", limit=10))
    project_ids = {item["source_project_id"] for item in result["matches"]}
    assert {"first", "second"} <= project_ids
    for match in result["matches"]:
        assert len(match["chunk_id"]) == 64
        assert len(match["source_sha256"]) == 64
        assert match["embedding_model"] == "test/embedding"
        assert match["chunk_version"] == CHUNK_VERSION
        assert match["index_fingerprint"]


def test_failed_rebuild_preserves_previous_collection(tmp_path: Path) -> None:
    service = make_service(tmp_path)
    vault = create_project(service)
    source = vault / "evidence.md"
    source.write_text("old reliable evidence", encoding="utf-8")
    asyncio.run(service.index("sample", reason="initial"))
    old_project = service.project("sample")
    stable = service.client.get_collection("project_sample")
    old_documents = stable.get(include=["documents"])["documents"]

    source.write_text("new evidence that cannot be embedded", encoding="utf-8")
    service.model.fail = True
    with pytest.raises(RuntimeError, match="embedding failed"):
        asyncio.run(service.index("sample", reason="failing rebuild"))

    stable = service.client.get_collection("project_sample")
    assert stable.get(include=["documents"])["documents"] == old_documents
    assert service.project("sample")["index_manifest"] == old_project["index_manifest"]


class DeleteFailingClient:
    def __init__(self, client) -> None:
        self.client = client

    def __getattr__(self, name):
        return getattr(self.client, name)

    def delete_collection(self, name: str) -> None:
        if name.startswith("deleted_"):
            raise RuntimeError("simulated Chroma delete failure")
        self.client.delete_collection(name)


def test_delete_failure_rolls_back_registry_and_collection(tmp_path: Path) -> None:
    service = make_service(tmp_path)
    vault = create_project(service)
    (vault / "evidence.md").write_text("keep me", encoding="utf-8")
    asyncio.run(service.index("sample", reason="initial"))
    service.client = DeleteFailingClient(service.client)

    with pytest.raises(RuntimeError, match="deletion failed"):
        asyncio.run(
            service.delete_project(
                "sample",
                ProjectDelete(confirm_name="Sample", trash_managed_files=True),
            )
        )

    assert service.project("sample")["id"] == "sample"
    assert vault.exists()
    assert service.client.get_collection("project_sample").count() > 0


def test_external_project_writes_to_managed_overlay(tmp_path: Path) -> None:
    service = make_service(tmp_path)
    external = service.safe(service.vault, "Projects/External")
    external.mkdir(parents=True)
    (external / "evidence.md").write_text("read only evidence", encoding="utf-8")
    service.save(
        [
            {
                "id": "external",
                "name": "External",
                "vault_path": "Projects/External",
                "source_paths": [],
                "indexed_at": None,
                "chunks": 0,
            }
        ]
    )
    asyncio.run(service.index("external", reason="legacy import"))

    async def capture_and_drain() -> dict:
        captured = await service.capture_progress(
            "external",
            ProgressCapture(
                session_id="session-1",
                phase="implementation",
                stage="review",
                title="Review",
                content="Completed safely",
                status="completed",
            ),
        )
        await service.wait_for_background_indexes()
        return captured

    result = asyncio.run(capture_and_drain())

    overlay = service.managed_root() / ".overlays" / "external"
    assert Path(result["note"]).is_file()
    assert result["index_pending"] is True
    assert result["index_status"] == "scheduled"
    assert overlay in Path(result["note"]).parents
    assert not (external / "Progress").exists()
    assert any(
        item["root_kind"] == "overlay"
        for item in service.project("external")["index_manifest"]["files"]
    )


def test_writeback_index_contract_syncs_research_and_coalesces_stage_notes(
    tmp_path: Path,
) -> None:
    service = make_service(tmp_path)
    vault = create_project(service)
    (vault / "evidence.md").write_text("initial evidence", encoding="utf-8")
    asyncio.run(service.index("sample", reason="initial"))

    async def exercise_contract() -> tuple[dict, dict, dict, dict]:
        research = await service.capture_research(
            "sample",
            ResearchCapture(
                query="verified source",
                sources=[
                    ResearchSource(
                        title="Primary source",
                        url="https://example.test/source",
                        content="synchronously searchable research evidence",
                    )
                ],
            ),
        )
        progress = await service.capture_progress(
            "sample",
            ProgressCapture(
                session_id="session-1",
                phase="implementation",
                stage="../outside/private-stage",
                title="Build",
                content="durable stage progress",
                status="completed",
            ),
        )
        handoff = await service.capture_handoff(
            "sample",
            CodexHandoff(
                goal="Finish the implementation",
                local_plan={"route": "codex"},
                local_response="Local analysis complete",
            ),
        )
        result = await service.capture_handoff_result(
            "sample",
            CodexHandoffResult(
                handoff_note=handoff["note"],
                worker_id="worker-1",
                attempt_id="attempt-1",
                success=True,
                summary="Implementation complete",
                output="Validated output",
                workspace_path="/workspace/sample",
            ),
        )
        duplicate = await service.capture_handoff_result(
            "sample",
            CodexHandoffResult(
                handoff_note=handoff["note"],
                worker_id="worker-1",
                attempt_id="attempt-1",
                success=True,
                summary="Implementation complete",
                output="Validated output",
                workspace_path="/workspace/sample",
            ),
        )
        await service.wait_for_background_indexes()
        assert duplicate["idempotent"] is True
        return research, progress, handoff, result

    research, progress, handoff, result = asyncio.run(exercise_contract())
    progress_note = Path(progress["note"])

    assert research["index_pending"] is False
    assert research["index_status"] == "rebuilt"
    assert progress["index_pending"] is True
    assert progress["index_status"] == "scheduled"
    assert "_agent-progress_" in progress_note.name
    assert "outside" not in progress_note.name
    assert "../outside/private-stage" in progress_note.read_text(encoding="utf-8")
    assert handoff["index_pending"] is True
    assert handoff["index_status"] == "coalesced"
    assert result["index_pending"] is True
    assert result["index_status"] == "coalesced"
    assert service.search("sample", Query(query="durable stage progress", limit=3))[
        "matches"
    ]
    assert "status: completed" in Path(result["note"]).read_text(encoding="utf-8")
    assert (
        Path(result["note"])
        .read_text(encoding="utf-8")
        .count("ai-workstation-codex-attempt:attempt-1")
        == 1
    )


def test_docx_xlsx_pptx_and_pdf_loaders(tmp_path: Path) -> None:
    service = make_service(tmp_path)
    fixtures = tmp_path / "fixtures"
    fixtures.mkdir()

    docx = fixtures / "sample.docx"
    with zipfile.ZipFile(docx, "w") as archive:
        archive.writestr(
            "word/document.xml",
            '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
            "<w:body><w:p><w:r><w:t>Document evidence</w:t></w:r></w:p></w:body>"
            "</w:document>",
        )
    assert "Document evidence" in service._extract_docx(docx)[0]

    from openpyxl import Workbook

    workbook = Workbook()
    workbook.active.append(["food", "value"])
    workbook.active.append(["apple", 42])
    xlsx = fixtures / "sample.xlsx"
    workbook.save(xlsx)
    assert "apple" in service._extract_xlsx(xlsx)[0]

    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Slide evidence"
    pptx = fixtures / "sample.pptx"
    presentation.save(pptx)
    assert "Slide evidence" in service._extract_pptx(pptx)[0]

    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=100, height=100)
    pdf = fixtures / "sample.pdf"
    with pdf.open("wb") as handle:
        writer.write(handle)
    assert service._extract_pdf(pdf)[0] == "## Page 1\n"
    assert {".pdf", ".docx", ".xlsx", ".xls", ".pptx"} == DOCUMENT_EXTENSIONS


def test_old_registry_shape_remains_readable(tmp_path: Path) -> None:
    service = make_service(tmp_path)
    vault = service.safe(service.vault, "Projects/Legacy")
    vault.mkdir(parents=True)
    (vault / "note.md").write_text("legacy", encoding="utf-8")
    service.save(
        [
            {
                "id": "legacy",
                "name": "Legacy",
                "vault_path": "Projects/Legacy",
                "source_paths": [],
                "indexed_at": None,
                "chunks": 0,
            }
        ]
    )

    assert service.project("legacy")["archived"] is False
    result = asyncio.run(service.index_if_stale("legacy", reason="migration"))
    assert result["rebuilt"] is True
    assert service.search("legacy", Query(query="legacy", limit=3))["matches"]


def test_orphan_cleanup_only_removes_unreferenced_old_uuid_directories(
    tmp_path: Path,
) -> None:
    service = make_service(tmp_path)
    vault = create_project(service)
    (vault / "evidence.md").write_text("active segment", encoding="utf-8")
    asyncio.run(service.index("sample", reason="initial"))
    chroma_root = service.data / "chroma"

    with sqlite3.connect(chroma_root / "chroma.sqlite3") as connection:
        referenced = {row[0] for row in connection.execute("SELECT id FROM segments")}
    referenced_directories = [
        chroma_root / identifier
        for identifier in referenced
        if (chroma_root / identifier).is_dir()
    ]
    assert referenced_directories

    orphan = chroma_root / str(uuid4())
    orphan.mkdir()
    (orphan / "index.bin").write_bytes(b"orphan")
    old = time.time() - 3600
    os.utime(orphan, (old, old))
    unrelated = chroma_root / "keep-me"
    unrelated.mkdir()

    service.orphan_grace_seconds = 0
    result = service.cleanup_orphan_segments()

    assert result["removed_directories"] >= 1
    assert not orphan.exists()
    assert unrelated.exists()
    assert all(path.exists() for path in referenced_directories)


def test_health_requires_completed_idle_sweep_and_detects_source_drift(
    tmp_path: Path,
) -> None:
    service = make_service(tmp_path)
    vault = create_project(service)
    source = vault / "evidence.md"
    source.write_text("stable evidence", encoding="utf-8")
    asyncio.run(service.index("sample", reason="initial"))
    service.watcher_state["running"] = True

    initializing = service.health_report()
    assert initializing["status"] == "degraded"
    assert initializing["sweep"]["last_completed_at"] is None

    service.sweep_state["last_completed_at"] = "2026-08-08T00:00:00+00:00"
    assert service.health_report()["status"] == "ok"

    service._background_index_pending["sample"] = {"writeback"}
    assert service.health_report()["status"] == "degraded"
    service._background_index_pending.clear()

    source.write_text("changed evidence with a new size", encoding="utf-8")
    changed = service.health_report()
    assert changed["status"] == "degraded"
    assert changed["index"]["stale_projects"] == 1
    assert "source files changed" in changed["index"]["projects"][0]["reasons"]


def test_health_report_does_not_expose_internal_exception_details(
    tmp_path: Path,
) -> None:
    service = make_service(tmp_path)
    secret = (
        "Traceback: private path C:" + r"\Users\owner\vault" + " and prompt contents"
    )
    service.watcher_state["last_error"] = secret
    service.sweep_state["last_error"] = secret
    service.background_index_state["last_error"] = secret
    service.cleanup_state["warnings"] = [secret]

    report = service.health_report()
    serialized = json.dumps(report)

    assert secret not in serialized
    assert report["watcher"]["last_error"] == "watcher operation failed"
    assert report["sweep"]["last_error"] == "consistency sweep failed"
    assert report["background_index"]["last_error"] == "background index failed"
    assert report["cleanup"]["warnings"] == [
        "cleanup operation reported a warning"
    ]


def test_safe_path_rejects_traversal_absolute_and_symlink_escape(
    tmp_path: Path,
) -> None:
    service = make_service(tmp_path)
    service.sources.mkdir()
    outside = tmp_path / "outside"
    outside.mkdir()

    for unsafe in (
        "../outside",
        "nested/../../outside",
        "/absolute/path",
        "C:" + r"\absolute\path",
        r"C:drive-relative",
        r"\\server\share\file",
        "nested//file",
    ):
        with pytest.raises(ValueError):
            service.safe(service.sources, unsafe)

    link = service.sources / "escape"
    try:
        link.symlink_to(outside, target_is_directory=True)
    except OSError as exc:  # pragma: no cover - unprivileged Windows runners
        pytest.skip(f"symbolic links unavailable: {exc}")
    with pytest.raises(ValueError):
        service.safe(service.sources, "escape/private.md")


@pytest.mark.parametrize(
    "source_path",
    [
        "../outside",
        "nested/../outside",
        "/absolute/path",
        "C:" + r"\absolute\path",
        r"\\server\share",
    ],
)
def test_project_model_rejects_unsafe_source_paths(source_path: str) -> None:
    with pytest.raises(ValueError):
        ProjectCreate(id="sample", name="Sample", source_paths=[source_path])


def test_project_creation_cannot_write_through_symlink_escape(tmp_path: Path) -> None:
    service = make_service(tmp_path)
    managed_root = service.managed_root()
    managed_root.mkdir(parents=True)
    outside = tmp_path / "outside"
    outside.mkdir()
    project_link = managed_root / "sample"
    try:
        project_link.symlink_to(outside, target_is_directory=True)
    except OSError as exc:  # pragma: no cover - unprivileged Windows runners
        pytest.skip(f"symbolic links unavailable: {exc}")

    with pytest.raises(ValueError):
        service.add(ProjectCreate(id="sample", name="Sample", source_paths=[]))
    assert not (outside / "00_Project.md").exists()


def test_managed_writeback_rejects_reparse_escape(tmp_path: Path) -> None:
    service = make_service(tmp_path)
    vault = create_project(service)
    outside = tmp_path / "outside"
    outside.mkdir()
    with pytest.raises(ValueError):
        service._atomic_write(outside / "escaped.md", "must stay managed")
    assert not (outside / "escaped.md").exists()

    handoffs = vault / "Handoffs"
    try:
        handoffs.symlink_to(outside, target_is_directory=True)
    except OSError as exc:  # pragma: no cover - unprivileged Windows runners
        pytest.skip(f"symbolic links unavailable: {exc}")

    with pytest.raises(ValueError):
        asyncio.run(
            service.capture_handoff(
                "sample",
                CodexHandoff(goal="Do not escape", local_plan={}),
            )
        )
    assert list(outside.iterdir()) == []


def test_handoff_result_is_limited_to_issued_managed_note(tmp_path: Path) -> None:
    service = make_service(tmp_path)
    vault = create_project(service)
    handoffs = vault / "Handoffs"
    issued = service._write_note(
        handoffs,
        "codex-handoff",
        "---\nstatus: pending\n---\n\n# Handoff\n",
    )

    def request(note: str) -> CodexHandoffResult:
        return CodexHandoffResult(
            handoff_note=note,
            worker_id="worker-1",
            success=True,
            summary="complete",
            output="validated",
            workspace_path="/workspace/sample",
        )

    outside = tmp_path / "outside"
    outside.mkdir()
    outside_note = outside / issued.name
    outside_note.write_text("secret", encoding="utf-8")
    for unsafe in (
        f"../{issued.name}",
        str(outside_note),
        "C:" + r"\foreign\2026-01-01_000000_000000_codex-handoff_abcdef.md",
    ):
        with pytest.raises(ValueError):
            asyncio.run(service.capture_handoff_result("sample", request(unsafe)))

    linked_note = handoffs / "2026-01-01_000000_000000_codex-handoff_abcdef.md"
    try:
        linked_note.symlink_to(outside_note)
    except OSError as exc:  # pragma: no cover - unprivileged Windows runners
        pytest.skip(f"symbolic links unavailable: {exc}")
    with pytest.raises(ValueError):
        asyncio.run(
            service.capture_handoff_result("sample", request(str(linked_note)))
        )

    async def record_valid_result() -> dict:
        result = await service.capture_handoff_result(
            "sample", request(issued.name)
        )
        await service.wait_for_background_indexes()
        return result

    result = asyncio.run(record_valid_result())
    assert result["status"] == "completed"
    assert "status: completed" in issued.read_text(encoding="utf-8")
