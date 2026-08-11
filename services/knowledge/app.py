from __future__ import annotations

import asyncio
import hashlib
import json
import os
import re
import shutil
import sqlite3
import stat
import threading
import xml.etree.ElementTree as ET
import zipfile
from collections import Counter
from contextlib import asynccontextmanager, suppress
from datetime import datetime, timezone
from pathlib import Path, PureWindowsPath
from typing import Any, Iterable
from uuid import UUID, uuid4

import chromadb
from fastapi import FastAPI, HTTPException
from fastapi import Query as FastAPIQuery
from pydantic import BaseModel, Field, field_validator
from sentence_transformers import SentenceTransformer
from watchfiles import awatch

VAULT, SOURCES, DATA = Path("/vault"), Path("/sources"), Path("/data")
MANAGED_PROJECT_ROOT = "Projects/AI Workbench"
REGISTRY = DATA / "projects.json"
PROJECT_ID_RE = re.compile(r"[a-z][a-z0-9-]{0,62}\Z")
HANDOFF_NOTE_RE = re.compile(
    r"\d{4}-\d{2}-\d{2}_\d{6}_\d{6}_codex-handoff_[0-9a-f]{6}\.md\Z"
)
NOTE_LABELS = frozenset({"agent-progress", "codex-handoff", "web-research"})

EMBEDDING_MODEL = os.getenv("KNOWLEDGE_EMBEDDING_MODEL", "BAAI/bge-m3")
MAX_TEXT_CHARS_PER_FILE = int(os.getenv("KNOWLEDGE_MAX_TEXT_CHARS_PER_FILE", "500000"))
CHUNK_SIZE = int(os.getenv("KNOWLEDGE_CHUNK_SIZE", "1200"))
CHUNK_OVERLAP = int(os.getenv("KNOWLEDGE_CHUNK_OVERLAP", "160"))
INDEX_SCHEMA_VERSION = 2
CHUNK_VERSION = f"characters-{CHUNK_SIZE}-overlap-{CHUNK_OVERLAP}-loaders-v2"
SWEEP_SECONDS = max(15, int(os.getenv("KNOWLEDGE_SWEEP_SECONDS", "300")))
WATCH_REFRESH_SECONDS = max(5, int(os.getenv("KNOWLEDGE_WATCH_REFRESH_SECONDS", "15")))
WATCH_POLL_MS = max(1000, int(os.getenv("KNOWLEDGE_WATCH_POLL_MS", "5000")))
ORPHAN_GRACE_SECONDS = max(0, int(os.getenv("KNOWLEDGE_ORPHAN_GRACE_SECONDS", "300")))

TEXT_EXTENSIONS = {
    ".md",
    ".txt",
    ".py",
    ".csv",
    ".json",
    ".jsonl",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
    ".cfg",
    ".conf",
    ".properties",
    ".xml",
    ".html",
    ".htm",
    ".css",
    ".scss",
    ".less",
    ".sql",
    ".tex",
    ".ipynb",
    ".swift",
    ".js",
    ".jsx",
    ".mjs",
    ".cjs",
    ".ts",
    ".tsx",
    ".vue",
    ".svelte",
    ".java",
    ".kt",
    ".kts",
    ".c",
    ".h",
    ".cc",
    ".cpp",
    ".hpp",
    ".cs",
    ".go",
    ".rs",
    ".rb",
    ".php",
    ".scala",
    ".sh",
    ".bash",
    ".zsh",
    ".fish",
    ".ps1",
    ".psm1",
    ".r",
    ".lua",
    ".dart",
    ".gradle",
    ".pbxproj",
    ".xcworkspacedata",
}
DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".xls", ".pptx"}
SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS | DOCUMENT_EXTENSIONS

IGNORED_DIRECTORY_NAMES = {
    ".git",
    ".hg",
    ".svn",
    ".cache",
    ".mypy_cache",
    ".pytest_cache",
    ".ruff_cache",
    ".tox",
    ".venv",
    ".idea",
    ".vscode",
    ".obsidian",
    ".trash",
    "__pycache__",
    "node_modules",
    "bower_components",
    "vendor",
    "pods",
    "deriveddata",
    "dist",
    "build",
    "out",
    "target",
    "coverage",
    ".next",
    ".nuxt",
    ".svelte-kit",
}


def utc_now() -> datetime:
    return datetime.now(timezone.utc)


def iso_now() -> str:
    return utc_now().isoformat()


def yaml_scalar(value: str) -> str:
    """JSON strings are valid YAML scalars and safely preserve punctuation/newlines."""
    return json.dumps(value, ensure_ascii=False)


class ProjectCreate(BaseModel):
    id: str = Field(pattern=r"^[a-z][a-z0-9-]{0,62}$")
    name: str = Field(min_length=1, max_length=120)
    source_paths: list[str] = Field(default_factory=list, max_length=10)

    @field_validator("name")
    @classmethod
    def normalize_name(cls, value: str) -> str:
        value = value.strip()
        if not value:
            raise ValueError("name must not be blank")
        return value

    @field_validator("source_paths")
    @classmethod
    def normalize_sources(cls, values: list[str]) -> list[str]:
        normalized: list[str] = []
        for value in values:
            value = value.strip().replace("\\", "/")
            windows_path = PureWindowsPath(value)
            if (
                not value
                or value.startswith("/")
                or windows_path.drive
                or windows_path.root
            ):
                raise ValueError("source paths must be relative subdirectories")
            parts = value.split("/")
            if any(part in {"", ".", ".."} for part in parts):
                raise ValueError("source paths must name a subdirectory")
            if value not in normalized:
                normalized.append(value)
        return normalized


class ProjectUpdate(BaseModel):
    name: str | None = Field(default=None, min_length=1, max_length=120)
    archived: bool | None = None

    @field_validator("name")
    @classmethod
    def normalize_name(cls, value: str | None) -> str | None:
        if value is None:
            return None
        value = value.strip()
        if not value:
            raise ValueError("name must not be blank")
        return value


class ProjectDelete(BaseModel):
    confirm_name: str = Field(min_length=1, max_length=120)
    trash_managed_files: bool = False


class Query(BaseModel):
    query: str = Field(min_length=1, max_length=2000)
    limit: int = Field(default=8, ge=1, le=20)


class ResearchSource(BaseModel):
    title: str = Field(min_length=1, max_length=500)
    url: str = Field(min_length=1, max_length=4000)
    content: str = Field(default="", max_length=20_000)


class ResearchCapture(BaseModel):
    query: str = Field(min_length=1, max_length=2000)
    sources: list[ResearchSource] = Field(min_length=1, max_length=10)


class ProgressCapture(BaseModel):
    session_id: str = Field(min_length=1, max_length=100)
    phase: str = Field(min_length=1, max_length=50)
    stage: str = Field(min_length=1, max_length=120)
    title: str = Field(min_length=1, max_length=200)
    content: str = Field(min_length=1, max_length=60_000)
    status: str = Field(min_length=1, max_length=50)


class CodexHandoff(BaseModel):
    goal: str = Field(min_length=1, max_length=20_000)
    local_plan: dict[str, Any] = Field(default_factory=dict)
    local_response: str = Field(default="", max_length=40_000)
    research_note: str | None = Field(default=None, max_length=4000)
    sources: list[ResearchSource] = Field(default_factory=list, max_length=10)
    workspace_id: str | None = Field(default=None, max_length=100)


class CodexHandoffResult(BaseModel):
    handoff_note: str = Field(min_length=1, max_length=4000)
    worker_id: str = Field(min_length=1, max_length=120)
    attempt_id: str | None = Field(
        default=None,
        min_length=1,
        max_length=120,
        pattern=r"^[A-Za-z0-9._:-]+$",
    )
    success: bool
    summary: str = Field(default="", max_length=20_000)
    output: str = Field(default="", max_length=60_000)
    workspace_path: str = Field(default="", max_length=2000)
    changed_files: list[str] = Field(default_factory=list, max_length=500)
    validation: list[str] = Field(default_factory=list, max_length=200)
    error: str | None = Field(default=None, max_length=20_000)


class KnowledgeService:
    def __init__(
        self,
        *,
        data: Path = DATA,
        vault: Path = VAULT,
        sources: Path = SOURCES,
        client: Any | None = None,
        model: Any | None = None,
        embedding_model: str = EMBEDDING_MODEL,
        embedding_revision: str | None = None,
        max_text_chars: int = MAX_TEXT_CHARS_PER_FILE,
        orphan_grace_seconds: int = ORPHAN_GRACE_SECONDS,
    ) -> None:
        self.data = Path(data)
        self.vault = Path(vault)
        self.sources = Path(sources)
        self.registry = self.data / "projects.json"
        self.max_text_chars = max_text_chars
        self.orphan_grace_seconds = max(0, orphan_grace_seconds)
        self.embedding_model = embedding_model
        self.data.mkdir(parents=True, exist_ok=True)
        self.state_lock = threading.RLock()
        self.model_lock = threading.Lock()
        self.index_lock = asyncio.Lock()
        self._background_index_tasks: dict[str, asyncio.Task[None]] = {}
        self._background_index_pending: dict[str, set[str]] = {}
        self._background_index_running: set[str] = set()
        self.background_index_state: dict[str, Any] = {
            "last_started_at": None,
            "last_completed_at": None,
            "last_error": None,
            "last_error_at": None,
        }
        self.client = client or chromadb.PersistentClient(
            path=str(self.data / "chroma")
        )
        self.embedding_revision = embedding_revision or self._model_revision()
        self.model = model or SentenceTransformer(
            self.embedding_model,
            cache_folder="/models",
            local_files_only=True,
        )
        self.watcher_state: dict[str, Any] = {
            "running": False,
            "roots": 0,
            "restarts": 0,
            "indexing_projects": [],
            "last_event_at": None,
            "last_error": None,
            "last_recovered_at": None,
        }
        self.sweep_state: dict[str, Any] = {
            "running": False,
            "last_started_at": None,
            "last_completed_at": None,
            "last_error": None,
        }
        self.cleanup_state: dict[str, Any] = {
            "last_run_at": None,
            "removed_directories": 0,
            "removed_bytes": 0,
            "skipped_recent": 0,
            "warnings": [],
        }
        self._recover_transient_collections()
        self.cleanup_orphan_segments()

    def _model_revision(self) -> str:
        cache_name = "models--" + self.embedding_model.replace("/", "--")
        reference = Path("/models") / cache_name / "refs" / "main"
        try:
            return reference.read_text(encoding="utf-8").strip() or "unknown"
        except OSError:
            return "unknown"

    @staticmethod
    def safe(root: Path, relative: str) -> Path:
        """Return a canonical descendant without accepting absolute or traversal paths.

        ``realpath`` is intentional: lexical normalization alone does not stop an
        existing symlink or Windows reparse point from escaping the configured root.
        The prefix guard is retained alongside ``commonpath`` so both humans and
        static analysis can prove that the returned path stays below the root.
        """
        if not isinstance(relative, str) or not relative or "\x00" in relative:
            raise ValueError("path must be a non-empty relative path")
        normalized = relative.replace("\\", "/")
        windows_path = PureWindowsPath(relative)
        if normalized.startswith("/") or windows_path.drive or windows_path.root:
            raise ValueError("absolute paths are not allowed")
        parts = normalized.split("/")
        if any(part in {"", ".", ".."} for part in parts):
            raise ValueError("path traversal is not allowed")

        resolved_root_text = os.path.realpath(
            os.path.normcase(os.path.abspath(os.fspath(root)))
        )
        candidate_text = os.path.realpath(
            os.path.normcase(os.path.join(resolved_root_text, *parts))
        )
        # Keep the recognized normalization -> prefix-check -> access sequence on
        # the exact value returned to callers.  Besides aiding static analysis,
        # this avoids validating a differently cased alias on Windows.
        if candidate_text.startswith(resolved_root_text):
            try:
                common = os.path.commonpath((resolved_root_text, candidate_text))
            except ValueError as exc:
                raise ValueError("path escapes configured root") from exc
            root_prefix = resolved_root_text.rstrip(os.sep) + os.sep
            if os.path.normcase(common) != resolved_root_text:
                raise ValueError("path escapes configured root")
            if candidate_text != resolved_root_text and not candidate_text.startswith(
                root_prefix
            ):
                raise ValueError("path escapes configured root")
            return Path(candidate_text)
        raise ValueError("path escapes configured root")

    @staticmethod
    def _validated_project_id(project_id: str) -> str:
        if PROJECT_ID_RE.fullmatch(project_id) is None:
            raise ValueError("invalid project id")
        return project_id

    @staticmethod
    def _is_reparse_point(path: Path) -> bool:
        try:
            metadata = path.lstat()
        except FileNotFoundError:
            return False
        if stat.S_ISLNK(metadata.st_mode):
            return True
        file_attributes = getattr(metadata, "st_file_attributes", 0)
        if file_attributes & getattr(stat, "FILE_ATTRIBUTE_REPARSE_POINT", 0):
            return True
        is_junction = getattr(path, "is_junction", None)
        return bool(is_junction and is_junction())

    @classmethod
    def _reject_reparse_components(cls, root: Path, candidate: Path) -> None:
        """Reject links/junctions below a trusted root before managed writes."""
        relative = candidate.relative_to(root)
        current = root
        for part in relative.parts:
            current = current / part
            if cls._is_reparse_point(current):
                raise ValueError("managed paths must not contain symbolic links")

    def _managed_path(self, candidate: Path) -> Path:
        """Canonicalize an absolute candidate and confine it to managed storage."""
        managed_root = self.managed_root()
        root_text = os.path.realpath(
            os.path.normcase(os.path.abspath(os.fspath(managed_root)))
        )
        lexical_text = os.path.abspath(os.path.normcase(os.fspath(candidate)))
        candidate_text = os.path.realpath(lexical_text)
        if candidate_text.startswith(root_text):
            try:
                common = os.path.commonpath((root_text, candidate_text))
            except ValueError as exc:
                raise ValueError("path escapes managed project storage") from exc
            root_prefix = root_text.rstrip(os.sep) + os.sep
            if os.path.normcase(common) != root_text:
                raise ValueError("path escapes managed project storage")
            if candidate_text != root_text and not candidate_text.startswith(root_prefix):
                raise ValueError("path escapes managed project storage")
            if lexical_text != root_text and not lexical_text.startswith(root_prefix):
                raise ValueError("path escapes managed project storage")

            resolved = Path(candidate_text)
            lexical = Path(lexical_text)
            try:
                lexical.relative_to(Path(root_text))
            except ValueError as exc:
                raise ValueError("path escapes managed project storage") from exc
            self._reject_reparse_components(Path(root_text), lexical)
            return resolved
        raise ValueError("path escapes managed project storage")

    def managed_root(self) -> Path:
        return self.safe(self.vault, MANAGED_PROJECT_ROOT)

    def _normalize_project(self, item: dict[str, Any]) -> dict[str, Any]:
        normalized = dict(item)
        normalized.setdefault("source_paths", [])
        normalized.setdefault("archived", False)
        normalized.setdefault("indexed_at", None)
        normalized.setdefault("chunks", 0)
        return normalized

    def _read_projects_unlocked(self) -> list[dict[str, Any]]:
        if not self.registry.exists():
            return []
        document = json.loads(self.registry.read_text(encoding="utf-8"))
        if not isinstance(document, list):
            raise ValueError("projects registry must contain a JSON array")
        return [self._normalize_project(item) for item in document]

    def projects(self) -> list[dict[str, Any]]:
        with self.state_lock:
            return self._read_projects_unlocked()

    def _save_unlocked(self, projects: list[dict[str, Any]]) -> None:
        temporary = self.registry.with_name(f".{self.registry.name}.{uuid4().hex}.tmp")
        try:
            with temporary.open("w", encoding="utf-8", newline="\n") as handle:
                json.dump(projects, handle, ensure_ascii=False, indent=2)
                handle.write("\n")
                handle.flush()
                os.fsync(handle.fileno())
            temporary.replace(self.registry)
        finally:
            with suppress(OSError):
                temporary.unlink()

    def save(self, projects: list[dict[str, Any]]) -> None:
        with self.state_lock:
            self._save_unlocked(projects)

    def _validated_source_roots(self, source_paths: Iterable[str]) -> list[Path]:
        roots: list[Path] = []
        for value in source_paths:
            value = value.strip().replace("\\", "/")
            windows_path = PureWindowsPath(value)
            if (
                not value
                or value.startswith("/")
                or windows_path.drive
                or windows_path.root
            ):
                raise ValueError("source paths must be relative subdirectories")
            if any(part in {"", ".", ".."} for part in value.split("/")):
                raise ValueError("source paths must name a subdirectory")
            root = self.safe(self.sources, value)
            if root == self.sources.resolve():
                raise ValueError("the complete source root cannot be indexed")
            if any(root == existing for existing in roots):
                continue
            if any(
                root in existing.parents or existing in root.parents
                for existing in roots
            ):
                raise ValueError("source paths must not overlap")
            roots.append(root)
        return roots

    def add(self, project: ProjectCreate) -> dict[str, Any]:
        self._validated_source_roots(project.source_paths)
        project_id = self._validated_project_id(project.id)
        with self.state_lock:
            projects = self._read_projects_unlocked()
            if any(item["id"] == project_id for item in projects):
                raise ValueError("project id already exists")
            managed_root = self.managed_root()
            managed_root.mkdir(parents=True, exist_ok=True)
            managed_root = self._managed_path(managed_root)
            vault = self.safe(managed_root, project_id)
            vault = self._managed_path(vault)
            vault.mkdir(parents=False, exist_ok=True)
            vault = self._managed_path(vault)
            note = self._managed_path(vault / "00_Project.md")
            if not note.exists():
                self._atomic_write(
                    note,
                    "---\n"
                    "type: project\n"
                    f"project: {project_id}\n"
                    "status: active\n"
                    "---\n\n"
                    f"# {project.name}\n",
                )
            item = project.model_dump() | {
                "id": project_id,
                "vault_path": f"{MANAGED_PROJECT_ROOT}/{project_id}",
                "indexed_at": None,
                "chunks": 0,
                "archived": False,
            }
            projects.append(item)
            self._save_unlocked(projects)
            return item

    def project(self, project_id: str) -> dict[str, Any]:
        with self.state_lock:
            project = next(
                (
                    item
                    for item in self._read_projects_unlocked()
                    if item["id"] == project_id
                ),
                None,
            )
        if project is None:
            raise KeyError(project_id)
        return project

    def _vault_path(self, project: dict[str, Any]) -> Path:
        return self.safe(self.vault, project["vault_path"])

    def _is_managed(self, vault: Path) -> bool:
        managed_root = self.managed_root()
        return managed_root == vault or managed_root in vault.parents

    def _overlay_path(self, project_id: str) -> Path:
        project_id = self._validated_project_id(project_id)
        return self.safe(self.managed_root(), f".overlays/{project_id}")

    def _writeback_path(self, project: dict[str, Any]) -> Path:
        vault = self._vault_path(project)
        candidate = (
            vault if self._is_managed(vault) else self._overlay_path(project["id"])
        )
        return self._managed_path(candidate)

    def _project_roots(self, project: dict[str, Any]) -> list[tuple[str, Path]]:
        vault = self._vault_path(project)
        roots: list[tuple[str, Path]] = [("obsidian", vault)]
        roots.extend(
            ("source", self.safe(self.sources, value))
            for value in project.get("source_paths", [])
        )
        if not self._is_managed(vault):
            roots.append(("overlay", self._overlay_path(project["id"])))
        unique: list[tuple[str, Path]] = []
        seen: set[Path] = set()
        for kind, root in roots:
            resolved = root.resolve()
            if resolved not in seen:
                seen.add(resolved)
                unique.append((kind, resolved))
        return unique

    def describe(self, project_id: str) -> dict[str, Any]:
        project = self.project(project_id)
        vault = self._vault_path(project)
        writeback = self._writeback_path(project)
        roots = [
            {"kind": kind, "path": self._display_path(root), "exists": root.exists()}
            for kind, root in self._project_roots(project)
        ]
        return {
            **self.public_project(project),
            "archived": bool(project.get("archived")),
            "managed": self._is_managed(vault),
            "writeback_path": self._display_path(writeback),
            "writeback_managed": True,
            "roots": roots,
        }

    @staticmethod
    def public_project(project: dict[str, Any]) -> dict[str, Any]:
        """Return management metadata without sending the full file manifest."""
        public = dict(project)
        manifest = dict(public.get("index_manifest") or {})
        if manifest:
            manifest.pop("files", None)
            public["index_manifest"] = manifest
        return public

    def update_project(self, project_id: str, request: ProjectUpdate) -> dict[str, Any]:
        with self.state_lock:
            projects = self._read_projects_unlocked()
            project = next(
                (item for item in projects if item["id"] == project_id), None
            )
            if project is None:
                raise KeyError(project_id)
            if request.name is not None:
                project["name"] = request.name.strip()
            if request.archived is not None:
                project["archived"] = request.archived
            self._save_unlocked(projects)
        return self.describe(project_id)

    async def delete_project(
        self, project_id: str, request: ProjectDelete
    ) -> dict[str, Any]:
        async with self.index_lock:
            return await asyncio.to_thread(
                self._delete_project_sync, project_id, request
            )

    def _delete_project_sync(
        self, project_id: str, request: ProjectDelete
    ) -> dict[str, Any]:
        stable_name = self._collection_name(project_id)
        tombstone_name = f"deleted_{project_id[:30]}_{uuid4().hex[:10]}"
        moved: list[tuple[Path, Path]] = []
        collection_renamed = False
        with self.state_lock:
            projects = self._read_projects_unlocked()
            project = next(
                (item for item in projects if item["id"] == project_id), None
            )
            if project is None:
                raise KeyError(project_id)
            if request.confirm_name != project["name"]:
                raise ValueError("confirmation name does not match project name")

            collection = self._get_collection(stable_name)
            try:
                if collection is not None:
                    collection.modify(name=tombstone_name)
                    collection_renamed = True
                if request.trash_managed_files:
                    vault = self._vault_path(project)
                    managed_path = (
                        vault
                        if self._is_managed(vault)
                        else self._overlay_path(project_id)
                    )
                    if managed_path.exists():
                        target = self._trash_target(project_id, managed_path)
                        managed_path.rename(target)
                        moved.append((managed_path, target))
                self._save_unlocked(
                    [item for item in projects if item["id"] != project_id]
                )
                if collection_renamed:
                    self.client.delete_collection(tombstone_name)
            except Exception as exc:
                rollback_errors: list[str] = []
                with suppress(Exception):
                    self._save_unlocked(projects)
                for original, target in reversed(moved):
                    try:
                        if target.exists() and not original.exists():
                            original.parent.mkdir(parents=True, exist_ok=True)
                            target.rename(original)
                    except (
                        Exception
                    ) as rollback_exc:  # pragma: no cover - rare I/O failure
                        rollback_errors.append(str(rollback_exc))
                if collection_renamed:
                    try:
                        renamed = self._get_collection(tombstone_name)
                        if renamed is not None:
                            renamed.modify(name=stable_name)
                    except Exception as rollback_exc:  # pragma: no cover
                        rollback_errors.append(str(rollback_exc))
                suffix = (
                    f"; rollback errors: {'; '.join(rollback_errors)}"
                    if rollback_errors
                    else ""
                )
                raise RuntimeError(
                    f"knowledge project deletion failed: {exc}{suffix}"
                ) from exc

        return {
            "deleted": True,
            "id": project_id,
            "name": project["name"],
            "index_deleted": True,
            "managed_files_trashed": bool(moved),
            "trashed_path": str(moved[0][1]) if moved else None,
            "trashed_paths": [str(target) for _, target in moved],
            "external_sources_deleted": False,
            "orphan_cleanup": self.cleanup_orphan_segments(),
        }

    def _trash_target(self, project_id: str, source: Path) -> Path:
        managed_root = self.managed_root()
        if managed_root != source and managed_root not in source.parents:
            raise ValueError("only workflow-managed files may be trashed")
        trash = managed_root / ".Trash"
        trash.mkdir(parents=True, exist_ok=True)
        stamp = datetime.now().strftime("%Y%m%d-%H%M%S-%f")
        label = "overlay" if ".overlays" in source.parts else "project"
        return trash / f"{project_id}-{label}-{stamp}"

    async def index_all(self) -> None:
        for project in self.projects():
            await self.index(project["id"], reason="index_all")

    async def index_stale(self, reason: str = "startup") -> list[dict[str, Any]]:
        results: list[dict[str, Any]] = []
        for project in self.projects():
            try:
                results.append(await self.index_if_stale(project["id"], reason=reason))
            except Exception as exc:
                self._record_index_error(project["id"], exc)
                results.append(
                    {
                        "project_id": project["id"],
                        "rebuilt": False,
                        "error": str(exc),
                    }
                )
        return results

    async def index_if_stale(self, project_id: str, *, reason: str) -> dict[str, Any]:
        async with self.index_lock:
            return await asyncio.to_thread(
                self._index_if_stale_sync, project_id, reason
            )

    def schedule_index(self, project_id: str, *, reason: str) -> dict[str, Any]:
        """Queue a coalesced rebuild without blocking a writeback request.

        Progress and handoff notes are durable as soon as their atomic write
        completes.  Re-embedding every PDF in a large project before returning
        made otherwise tiny writebacks block workers for minutes, so those
        writebacks use this queue.  The filesystem watcher and consistency sweep
        remain independent fallbacks if the in-process task is interrupted.
        """
        project = self.project(project_id)
        loop = asyncio.get_running_loop()
        with self.state_lock:
            self._background_index_pending.setdefault(project_id, set()).add(reason)
            task = self._background_index_tasks.get(project_id)
            if task is None or task.done():
                task = loop.create_task(
                    self._background_index_worker(project_id),
                    name=f"knowledge-index-{project_id}",
                )
                self._background_index_tasks[project_id] = task
                status = "scheduled"
            else:
                status = "coalesced"
        return {
            "project_id": project_id,
            "chunks": int(project.get("chunks") or 0),
            "indexed_at": project.get("indexed_at"),
            "rebuilt": False,
            "reason": reason,
            "index_pending": True,
            "index_status": status,
        }

    async def _background_index_worker(self, project_id: str) -> None:
        current_task = asyncio.current_task()
        try:
            # Give several stage events from one workflow turn a small window to
            # coalesce into one source scan and one atomic collection swap.
            await asyncio.sleep(0.25)
            while True:
                with self.state_lock:
                    reasons = self._background_index_pending.pop(project_id, set())
                    if not reasons:
                        break
                    self._background_index_running.add(project_id)
                    self.background_index_state.update(
                        {"last_started_at": iso_now(), "last_error": None}
                    )
                try:
                    await self.index_if_stale(
                        project_id,
                        reason="background writeback: " + ", ".join(sorted(reasons)),
                    )
                except asyncio.CancelledError:
                    raise
                except Exception as exc:
                    self._record_index_error(project_id, exc)
                    with self.state_lock:
                        self.background_index_state.update(
                            {
                                "last_error": str(exc)[:4000],
                                "last_error_at": iso_now(),
                            }
                        )
                finally:
                    with self.state_lock:
                        self._background_index_running.discard(project_id)
                        self.background_index_state["last_completed_at"] = iso_now()
                # A note written during the rebuild leaves another pending reason.
                # Recheck it after yielding so no event is lost at the scan boundary.
                await asyncio.sleep(0)
        finally:
            with self.state_lock:
                self._background_index_running.discard(project_id)
                if self._background_index_tasks.get(project_id) is current_task:
                    self._background_index_tasks.pop(project_id, None)

    async def wait_for_background_indexes(self) -> None:
        """Drain queued rebuilds (used by shutdown and deterministic tests)."""
        while True:
            with self.state_lock:
                tasks = [
                    task
                    for task in self._background_index_tasks.values()
                    if not task.done()
                ]
            if not tasks:
                return
            await asyncio.gather(*tasks, return_exceptions=True)

    async def cancel_background_indexes(self) -> None:
        with self.state_lock:
            tasks = [
                task
                for task in self._background_index_tasks.values()
                if not task.done()
            ]
        for task in tasks:
            task.cancel()
        if tasks:
            await asyncio.gather(*tasks, return_exceptions=True)

    def _index_if_stale_sync(self, project_id: str, reason: str) -> dict[str, Any]:
        project = self.project(project_id)
        inventory = self._scan_project(project, extract=False)
        reasons = self._stale_reasons(project, inventory)
        if not reasons:
            return {
                "project_id": project_id,
                "chunks": int(project.get("chunks") or 0),
                "indexed_at": project.get("indexed_at"),
                "rebuilt": False,
                "reason": "current",
                "fingerprint": inventory["manifest"]["fingerprint"],
            }
        return self._rebuild_sync(
            project_id,
            reason=f"{reason}: {', '.join(reasons)}",
        )

    async def index(self, project_id: str, *, reason: str = "manual") -> dict[str, Any]:
        async with self.index_lock:
            return await asyncio.to_thread(self._rebuild_sync, project_id, reason)

    def _rebuild_sync(self, project_id: str, reason: str) -> dict[str, Any]:
        project = self.project(project_id)
        scan = self._scan_project(project, extract=True)
        stats = scan["stats"]
        if stats["failed_files"]:
            failures = "; ".join(stats["failures"][:3])
            raise RuntimeError(
                f"index scan failed for {stats['failed_files']} file(s): {failures}"
            )

        stable_name = self._collection_name(project_id)
        build_name = f"build_{project_id[:30]}_{uuid4().hex[:10]}"
        backup_name = f"backup_{project_id[:30]}_{uuid4().hex[:10]}"
        metadata = {
            "project_id": project_id,
            "index_schema": INDEX_SCHEMA_VERSION,
            "chunk_version": CHUNK_VERSION,
            "embedding_model": self.embedding_model,
            "embedding_revision": self.embedding_revision,
            "fingerprint": scan["manifest"]["fingerprint"],
            "built_at": iso_now(),
        }
        temporary = None
        try:
            temporary = self.client.create_collection(build_name, metadata=metadata)
            documents = scan["documents"]
            if documents:
                # Do not hold the inference lock for an entire large project.
                # Search queries can run between batches while the old stable
                # collection continues serving traffic.
                for start in range(0, len(documents), 64):
                    batch = documents[start : start + 64]
                    with self.model_lock:
                        vectors = self.model.encode(
                            [item["text"] for item in batch],
                            normalize_embeddings=True,
                            show_progress_bar=False,
                            batch_size=32,
                        ).tolist()
                    temporary.add(
                        ids=[item["id"] for item in batch],
                        documents=[item["text"] for item in batch],
                        embeddings=vectors,
                        metadatas=[item["metadata"] for item in batch],
                    )
                if temporary.count() != len(documents):
                    raise RuntimeError("temporary collection count mismatch")
        except Exception:
            with suppress(Exception):
                self.client.delete_collection(build_name)
            raise

        backup_created = False
        switched = False
        with self.state_lock:
            projects = self._read_projects_unlocked()
            latest = next((item for item in projects if item["id"] == project_id), None)
            if latest is None:
                with suppress(Exception):
                    self.client.delete_collection(build_name)
                raise KeyError(project_id)
            old = self._get_collection(stable_name)
            try:
                if old is not None:
                    old.modify(name=backup_name)
                    backup_created = True
                temporary = self._get_collection(build_name)
                if temporary is None:
                    raise RuntimeError("temporary collection disappeared before switch")
                temporary.modify(name=stable_name)
                switched = True

                now = iso_now()
                stats = {
                    **stats,
                    "last_reason": reason,
                    "last_error": None,
                    "last_error_at": None,
                    "completed_at": now,
                }
                latest.update(
                    {
                        "indexed_at": now,
                        "chunks": len(scan["documents"]),
                        "index_manifest": scan["manifest"],
                        "index_stats": stats,
                    }
                )
                self._save_unlocked(projects)
            except Exception:
                if switched:
                    with suppress(Exception):
                        active = self._get_collection(stable_name)
                        if active is not None:
                            active.modify(name=build_name)
                if backup_created:
                    with suppress(Exception):
                        backup = self._get_collection(backup_name)
                        if backup is not None:
                            backup.modify(name=stable_name)
                with suppress(Exception):
                    self.client.delete_collection(build_name)
                raise

            cleanup_error: str | None = None
            if backup_created:
                try:
                    self.client.delete_collection(backup_name)
                except Exception as exc:  # old index no longer serves queries
                    cleanup_error = str(exc)
                    latest_stats = dict(latest.get("index_stats") or {})
                    latest_stats["cleanup_warning"] = cleanup_error
                    latest["index_stats"] = latest_stats
                    self._save_unlocked(projects)

        return {
            "project_id": project_id,
            "chunks": len(scan["documents"]),
            "indexed_at": latest["indexed_at"],
            "rebuilt": True,
            "reason": reason,
            "fingerprint": scan["manifest"]["fingerprint"],
            "index_stats": latest["index_stats"],
            "orphan_cleanup": self.cleanup_orphan_segments(),
        }

    def _stale_reasons(
        self, project: dict[str, Any], inventory: dict[str, Any]
    ) -> list[str]:
        reasons: list[str] = []
        manifest = project.get("index_manifest") or {}
        if not manifest:
            reasons.append("manifest missing")
        if manifest.get("fingerprint") != inventory["manifest"]["fingerprint"]:
            reasons.append("source fingerprint changed")
        if manifest.get("embedding_model") != self.embedding_model:
            reasons.append("embedding model changed")
        if manifest.get("embedding_revision") != self.embedding_revision:
            reasons.append("embedding revision changed")
        if manifest.get("chunk_version") != CHUNK_VERSION:
            reasons.append("chunk version changed")
        if manifest.get("max_text_chars") != self.max_text_chars:
            reasons.append("text limit changed")
        collection = self._get_collection(self._collection_name(project["id"]))
        if collection is None:
            reasons.append("collection missing")
        else:
            try:
                if collection.count() != int(project.get("chunks") or 0):
                    reasons.append("collection count mismatch")
            except Exception:
                reasons.append("collection unreadable")
        return list(dict.fromkeys(reasons))

    def _record_index_error(self, project_id: str, error: Exception) -> None:
        with self.state_lock:
            projects = self._read_projects_unlocked()
            project = next(
                (item for item in projects if item["id"] == project_id), None
            )
            if project is None:
                return
            stats = dict(project.get("index_stats") or {})
            stats.update({"last_error": str(error)[:4000], "last_error_at": iso_now()})
            project["index_stats"] = stats
            self._save_unlocked(projects)

    def _scan_project(
        self, project: dict[str, Any], *, extract: bool
    ) -> dict[str, Any]:
        stats: dict[str, Any] = {
            "roots": 0,
            "missing_roots": [],
            "scanned_files": 0,
            "supported_files": 0,
            "indexed_files": 0,
            "empty_files": 0,
            "unsupported_files": 0,
            "unsupported_extensions": {},
            "truncated_files": [],
            "failed_files": 0,
            "failures": [],
            "chunks": 0,
        }
        manifest_files: list[dict[str, Any]] = []
        unsupported_inventory: list[dict[str, Any]] = []
        documents: list[dict[str, Any]] = []
        unsupported = Counter()

        for kind, root in self._project_roots(project):
            stats["roots"] += 1
            if not root.exists():
                # An overlay is created lazily on the first writeback.  Its absence
                # is normal for an external read-only project and must not make the
                # source project look unhealthy.
                if kind != "overlay":
                    stats["missing_roots"].append(self._display_path(root))
                continue
            for path in self._iter_files(root):
                stats["scanned_files"] += 1
                extension = path.suffix.lower()
                try:
                    stat_before = path.stat()
                except OSError as exc:
                    stats["failed_files"] += 1
                    stats["failures"].append(f"{path}: {exc}")
                    continue
                if extension not in SUPPORTED_EXTENSIONS:
                    unsupported[extension or "<none>"] += 1
                    unsupported_inventory.append(
                        {
                            "source": self._display_path(path),
                            "size": stat_before.st_size,
                            "mtime_ns": stat_before.st_mtime_ns,
                        }
                    )
                    continue

                stats["supported_files"] += 1
                entry: dict[str, Any] = {
                    "source": self._display_path(path),
                    "root_kind": kind,
                    "extension": extension,
                    "size": stat_before.st_size,
                    "mtime_ns": stat_before.st_mtime_ns,
                    "sha256": None,
                    "chunks": 0,
                    "truncated": False,
                    "status": "ready",
                }
                try:
                    entry["sha256"] = self._file_sha256(path)
                    if extract:
                        text, observed_chars, complete = self._extract_file(path)
                        stat_after = path.stat()
                        if (
                            stat_after.st_size != stat_before.st_size
                            or stat_after.st_mtime_ns != stat_before.st_mtime_ns
                        ):
                            raise RuntimeError(
                                "file changed while it was being indexed"
                            )
                        normalized = text.strip()
                        truncated = (
                            not complete or len(normalized) > self.max_text_chars
                        )
                        if truncated:
                            stats["truncated_files"].append(
                                {
                                    "source": self._display_path(path),
                                    "observed_chars": observed_chars,
                                    "indexed_chars": min(
                                        len(normalized), self.max_text_chars
                                    ),
                                    "complete_count": complete,
                                }
                            )
                            normalized = normalized[: self.max_text_chars]
                        entry["truncated"] = truncated
                        entry["observed_chars"] = observed_chars
                        chunks = self.chunks(normalized)
                        entry["chunks"] = len(chunks)
                        if chunks:
                            stats["indexed_files"] += 1
                            for index, chunk in enumerate(chunks):
                                digest = hashlib.sha256(
                                    (
                                        f"{self._display_path(path)}|{CHUNK_VERSION}|"
                                        f"{index}|{chunk}"
                                    ).encode("utf-8")
                                ).hexdigest()
                                documents.append(
                                    {
                                        "id": digest,
                                        "text": chunk,
                                        "metadata": {
                                            "source": self._display_path(path),
                                            "chunk": index,
                                            "source_sha256": entry["sha256"],
                                            "root_kind": kind,
                                        },
                                    }
                                )
                        else:
                            stats["empty_files"] += 1
                            entry["status"] = "empty"
                except Exception as exc:
                    entry["status"] = "failed"
                    entry["error"] = str(exc)[:1000]
                    stats["failed_files"] += 1
                    stats["failures"].append(f"{self._display_path(path)}: {exc}")
                manifest_files.append(entry)

        stats["unsupported_files"] = sum(unsupported.values())
        stats["unsupported_extensions"] = dict(sorted(unsupported.items()))
        stats["chunks"] = len(documents) if extract else int(project.get("chunks") or 0)
        manifest_files.sort(key=lambda item: item["source"])
        unsupported_inventory.sort(key=lambda item: item["source"])
        fingerprint_payload = {
            "index_schema": INDEX_SCHEMA_VERSION,
            "chunk_version": CHUNK_VERSION,
            "max_text_chars": self.max_text_chars,
            "embedding_model": self.embedding_model,
            "embedding_revision": self.embedding_revision,
            "files": [
                {
                    "source": item["source"],
                    "size": item["size"],
                    "sha256": item["sha256"],
                }
                for item in manifest_files
            ],
            "unsupported_inventory": unsupported_inventory,
            "missing_roots": stats["missing_roots"],
        }
        fingerprint = hashlib.sha256(
            json.dumps(
                fingerprint_payload,
                ensure_ascii=False,
                sort_keys=True,
                separators=(",", ":"),
            ).encode("utf-8")
        ).hexdigest()
        manifest = {
            "schema_version": INDEX_SCHEMA_VERSION,
            "fingerprint": fingerprint,
            "embedding_model": self.embedding_model,
            "embedding_revision": self.embedding_revision,
            "chunk_version": CHUNK_VERSION,
            "max_text_chars": self.max_text_chars,
            "files": manifest_files,
            "generated_at": iso_now(),
        }
        return {"manifest": manifest, "stats": stats, "documents": documents}

    def _iter_files(self, root: Path) -> Iterable[Path]:
        if root.is_file():
            yield root
            return
        for directory, directories, files in os.walk(root, followlinks=False):
            directories[:] = sorted(
                name
                for name in directories
                if not self._ignore_directory(name)
                and not (Path(directory) / name).is_symlink()
            )
            for name in sorted(files):
                path = Path(directory) / name
                if path.is_symlink() or name.startswith(".~") or name.endswith("~"):
                    continue
                yield path

    @staticmethod
    def _ignore_directory(name: str) -> bool:
        lowered = name.lower()
        if lowered in IGNORED_DIRECTORY_NAMES:
            return True
        return name.startswith(".") and lowered != ".github"

    @staticmethod
    def _file_sha256(path: Path) -> str:
        digest = hashlib.sha256()
        with path.open("rb") as handle:
            for block in iter(lambda: handle.read(1024 * 1024), b""):
                digest.update(block)
        return digest.hexdigest()

    def _extract_file(self, path: Path) -> tuple[str, int, bool]:
        extension = path.suffix.lower()
        if extension in TEXT_EXTENSIONS:
            raw = path.read_bytes()
            text = self._decode_text(raw)
            if extension == ".ipynb":
                text = self._extract_notebook(text)
            return text, len(text), True
        if extension == ".pdf":
            return self._extract_pdf(path)
        if extension == ".docx":
            return self._extract_docx(path)
        if extension == ".xlsx":
            return self._extract_xlsx(path)
        if extension == ".xls":
            return self._extract_xls(path)
        if extension == ".pptx":
            return self._extract_pptx(path)
        raise ValueError(f"unsupported extension: {extension}")

    @staticmethod
    def _decode_text(raw: bytes) -> str:
        for encoding in ("utf-8-sig", "utf-16", "gb18030"):
            try:
                return raw.decode(encoding)
            except UnicodeDecodeError:
                continue
        raise UnicodeDecodeError(
            "utf-8", raw, 0, min(1, len(raw)), "unsupported text encoding"
        )

    @staticmethod
    def _extract_notebook(content: str) -> str:
        document = json.loads(content)
        lines: list[str] = []
        for cell in document.get("cells", []):
            cell_type = cell.get("cell_type", "cell")
            source = cell.get("source", [])
            text = "".join(source) if isinstance(source, list) else str(source)
            lines.extend((f"## {cell_type}", text, ""))
        return "\n".join(lines)

    def _bounded_join(self, values: Iterable[str]) -> tuple[str, int, bool]:
        parts: list[str] = []
        observed = 0
        hard_limit = self.max_text_chars + 1
        complete = True
        for value in values:
            if not value:
                continue
            value = str(value)
            remaining = hard_limit - observed
            if remaining <= 0:
                complete = False
                break
            segment = value[:remaining]
            parts.append(segment)
            observed += len(segment)
            if len(segment) < len(value):
                complete = False
                break
        return "\n".join(parts), observed, complete

    def _extract_pdf(self, path: Path) -> tuple[str, int, bool]:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        return self._bounded_join(
            f"## Page {index + 1}\n{page.extract_text() or ''}"
            for index, page in enumerate(reader.pages)
        )

    def _extract_docx(self, path: Path) -> tuple[str, int, bool]:
        with zipfile.ZipFile(path) as archive:
            names = [
                name
                for name in archive.namelist()
                if name == "word/document.xml"
                or name.startswith("word/header")
                or name.startswith("word/footer")
                or name in {"word/footnotes.xml", "word/endnotes.xml"}
            ]
            values: list[str] = []
            for name in sorted(names, key=lambda value: value != "word/document.xml"):
                root = ET.fromstring(archive.read(name))
                values.append(
                    " ".join(part for part in root.itertext() if part.strip())
                )
        return self._bounded_join(values)

    def _extract_xlsx(self, path: Path) -> tuple[str, int, bool]:
        from openpyxl import load_workbook

        workbook = load_workbook(path, read_only=True, data_only=True)

        def rows() -> Iterable[str]:
            try:
                for sheet in workbook.worksheets:
                    yield f"## Sheet: {sheet.title}"
                    for row in sheet.iter_rows(values_only=True):
                        values = ["" if value is None else str(value) for value in row]
                        if any(values):
                            yield "\t".join(values)
            finally:
                workbook.close()

        return self._bounded_join(rows())

    def _extract_xls(self, path: Path) -> tuple[str, int, bool]:
        import xlrd

        workbook = xlrd.open_workbook(str(path), on_demand=True)

        def rows() -> Iterable[str]:
            try:
                for sheet in workbook.sheets():
                    yield f"## Sheet: {sheet.name}"
                    for row_index in range(sheet.nrows):
                        values = [str(value) for value in sheet.row_values(row_index)]
                        if any(values):
                            yield "\t".join(values)
            finally:
                workbook.release_resources()

        return self._bounded_join(rows())

    def _extract_pptx(self, path: Path) -> tuple[str, int, bool]:
        from pptx import Presentation

        presentation = Presentation(str(path))

        def slides() -> Iterable[str]:
            for index, slide in enumerate(presentation.slides):
                yield f"## Slide {index + 1}"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        yield shape.text

        return self._bounded_join(slides())

    @staticmethod
    def chunks(
        content: str,
        size: int = CHUNK_SIZE,
        overlap: int = CHUNK_OVERLAP,
    ) -> list[str]:
        content = content.strip()
        if not content:
            return []
        step = size - overlap
        if step <= 0:
            raise ValueError("chunk overlap must be smaller than chunk size")
        return [content[start : start + size] for start in range(0, len(content), step)]

    def documents(self, project: dict[str, Any]) -> list[dict[str, Any]]:
        scan = self._scan_project(project, extract=True)
        if scan["stats"]["failed_files"]:
            raise RuntimeError("one or more files could not be extracted")
        return scan["documents"]

    def search(self, project_id: str, request: Query) -> dict[str, Any]:
        project = self.project(project_id)
        with self.model_lock:
            vector = self.model.encode(
                [request.query],
                normalize_embeddings=True,
                show_progress_bar=False,
            ).tolist()
        with self.state_lock:
            collection = self._get_collection(self._collection_name(project_id))
            if collection is None:
                raise KeyError(project_id)
            count = collection.count()
            if count == 0:
                return {"project_id": project_id, "query": request.query, "matches": []}
            result = collection.query(
                query_embeddings=vector,
                n_results=min(request.limit, count),
                include=["documents", "metadatas", "distances"],
            )
            index_metadata = dict(collection.metadata or {})
        return {
            "project_id": project_id,
            "query": request.query,
            "matches": [
                {
                    "chunk_id": chunk_id,
                    "text": text,
                    "source": meta.get("source"),
                    "chunk": meta.get("chunk"),
                    "source_sha256": meta.get("source_sha256"),
                    "root_kind": meta.get("root_kind"),
                    "source_project_id": project_id,
                    "source_project_name": project.get("name"),
                    "index_fingerprint": index_metadata.get("fingerprint"),
                    "embedding_model": index_metadata.get("embedding_model"),
                    "chunk_version": index_metadata.get("chunk_version"),
                    "indexed_at": index_metadata.get("built_at"),
                    "distance": distance,
                }
                for chunk_id, text, meta, distance in zip(
                    result["ids"][0],
                    result["documents"][0],
                    result["metadatas"][0],
                    result["distances"][0],
                )
            ],
        }

    def search_all(self, request: Query) -> dict[str, Any]:
        """Search every active project and return globally ranked evidence."""

        with self.model_lock:
            vector = self.model.encode(
                [request.query],
                normalize_embeddings=True,
                show_progress_bar=False,
            ).tolist()
        matches: list[dict[str, Any]] = []
        with self.state_lock:
            projects = [
                project
                for project in self._read_projects_unlocked()
                if not project.get("archived")
            ]
            for project in projects:
                collection = self._get_collection(
                    self._collection_name(project["id"])
                )
                if collection is None or collection.count() == 0:
                    continue
                result = collection.query(
                    query_embeddings=vector,
                    n_results=min(request.limit, collection.count()),
                    include=["documents", "metadatas", "distances"],
                )
                index_metadata = dict(collection.metadata or {})
                matches.extend(
                    {
                        "chunk_id": chunk_id,
                        "text": text,
                        "source": meta.get("source"),
                        "chunk": meta.get("chunk"),
                        "source_sha256": meta.get("source_sha256"),
                        "root_kind": meta.get("root_kind"),
                        "source_project_id": project["id"],
                        "source_project_name": project.get("name"),
                        "index_fingerprint": index_metadata.get("fingerprint"),
                        "embedding_model": index_metadata.get("embedding_model"),
                        "chunk_version": index_metadata.get("chunk_version"),
                        "indexed_at": index_metadata.get("built_at"),
                        "distance": distance,
                    }
                    for chunk_id, text, meta, distance in zip(
                        result["ids"][0],
                        result["documents"][0],
                        result["metadatas"][0],
                        result["distances"][0],
                    )
                )
        matches.sort(key=lambda item: float(item.get("distance", float("inf"))))
        unique: list[dict[str, Any]] = []
        seen: set[tuple[str, str]] = set()
        for match in matches:
            key = (
                str(match.get("source_project_id") or ""),
                str(match.get("chunk_id") or ""),
            )
            if key in seen:
                continue
            seen.add(key)
            unique.append(match)
            if len(unique) >= request.limit:
                break
        return {"query": request.query, "matches": unique}

    async def capture_research(
        self, project_id: str, request: ResearchCapture
    ) -> dict[str, Any]:
        project = self.project(project_id)
        directory = self._writeback_path(project) / "Research"
        lines = [
            "---",
            "type: web-research",
            f"project: {project_id}",
            f"query: {yaml_scalar(request.query)}",
            f"captured_at: {iso_now()}",
            "---",
            "",
            f"# Web research: {request.query}",
            "",
        ]
        for source in request.sources:
            lines += [
                f"## {source.title}",
                "",
                f"Source: {source.url}",
                "",
                source.content.strip() or "_No summary returned by search provider._",
                "",
            ]
        note = self._write_note(directory, "web-research", "\n".join(lines))
        indexed = await self.index_if_stale(project_id, reason="research writeback")
        return {
            "project_id": project_id,
            "note": str(note),
            "sources": len(request.sources),
            **indexed,
            "index_pending": False,
            "index_status": "rebuilt" if indexed.get("rebuilt") else "current",
        }

    async def capture_progress(
        self, project_id: str, request: ProgressCapture
    ) -> dict[str, Any]:
        project = self.project(project_id)
        directory = self._writeback_path(project) / "Progress"
        lines = [
            "---",
            "type: agent-progress",
            f"project: {project_id}",
            f"agent_session: {yaml_scalar(request.session_id)}",
            f"phase: {yaml_scalar(request.phase)}",
            f"stage: {yaml_scalar(request.stage)}",
            f"status: {yaml_scalar(request.status)}",
            f"captured_at: {iso_now()}",
            "---",
            "",
            f"# {request.title}",
            "",
            request.content.strip(),
            "",
        ]
        # Stage remains in the document metadata; filenames use a server-owned
        # category so request data never becomes a filesystem path component.
        note = self._write_note(directory, "agent-progress", "\n".join(lines))
        indexed = self.schedule_index(project_id, reason="progress writeback")
        return {"project_id": project_id, "note": str(note), **indexed}

    async def capture_handoff(
        self, project_id: str, request: CodexHandoff
    ) -> dict[str, Any]:
        project = self.project(project_id)
        directory = self._writeback_path(project) / "Handoffs"
        lines = [
            "---",
            "type: codex-handoff",
            f"project: {project_id}",
            "status: pending",
            f"created_at: {iso_now()}",
            f"workspace_id: {yaml_scalar(request.workspace_id or '')}",
            "---",
            "",
            "# Codex implementation handoff",
            "",
            "## Goal",
            "",
            request.goal,
            "",
            "## Local planner decision",
            "",
            "```json",
            json.dumps(request.local_plan, ensure_ascii=False, indent=2),
            "```",
            "",
            "## Local-model analysis",
            "",
            request.local_response or "_No final local-model analysis was returned._",
            "",
        ]
        if request.research_note:
            lines += ["## Research note", "", request.research_note, ""]
        if request.sources:
            lines += ["## Verified web sources", ""]
            for source in request.sources:
                lines += [f"- [{source.title}]({source.url})", ""]
        lines += [
            "## Completion contract",
            "",
            "- Inspect the listed project evidence before making changes.",
            "- Keep changes within the named project/workspace scope.",
            "- Record implementation files, validation performed, and remaining risks in this note.",
            "- Change `status` from `pending` to `completed` only after validation.",
            "",
        ]
        note = self._write_note(directory, "codex-handoff", "\n".join(lines))
        indexed = self.schedule_index(project_id, reason="handoff writeback")
        return {"project_id": project_id, "note": str(note), **indexed}

    async def capture_handoff_result(
        self, project_id: str, request: CodexHandoffResult
    ) -> dict[str, Any]:
        project = self.project(project_id)
        directory = self._managed_path(self._writeback_path(project) / "Handoffs")
        note = self._handoff_note_path(directory, request.handoff_note)
        # Re-enumerate immediately before reading so the request is used only as
        # an equality selector; the filesystem path itself comes from the server.
        note = self._handoff_note_path(directory, request.handoff_note)
        content = note.read_text(encoding="utf-8")
        attempt_marker = (
            f"<!-- ai-workstation-codex-attempt:{request.attempt_id} -->"
            if request.attempt_id
            else None
        )
        if attempt_marker and attempt_marker in content:
            return {
                "project_id": project_id,
                "note": str(note),
                "status": "already_recorded",
                "rebuilt": False,
                "index_pending": False,
                "index_status": "already_recorded",
                "idempotent": True,
            }
        final_status = "completed" if request.success else "failed"
        content = content.replace("status: pending", f"status: {final_status}", 1)
        lines = [
            "",
            *([attempt_marker, ""] if attempt_marker else []),
            "## Codex execution result",
            "",
            f"- Status: `{final_status}`",
            f"- Worker: `{request.worker_id}`",
            f"- Completed: `{iso_now()}`",
            f"- Workspace: `{request.workspace_path}`",
            "",
            "### Summary",
            "",
            request.summary or "_No summary returned._",
            "",
        ]
        if request.changed_files:
            lines += (
                ["### Changed files", ""]
                + [f"- `{path}`" for path in request.changed_files]
                + [""]
            )
        if request.validation:
            lines += (
                ["### Validation", ""]
                + [f"- {item}" for item in request.validation]
                + [""]
            )
        if request.error:
            lines += ["### Error", "", request.error, ""]
        lines += [
            "### Codex output",
            "",
            request.output or "_No output returned._",
            "",
        ]
        note = self._handoff_note_path(directory, request.handoff_note)
        self._atomic_write(note, content.rstrip() + "\n" + "\n".join(lines))
        indexed = self.schedule_index(project_id, reason="handoff result")
        return {
            "project_id": project_id,
            "note": str(note),
            "status": final_status,
            **indexed,
        }

    def _handoff_note_path(self, directory: Path, supplied: str) -> Path:
        """Select an issued handoff from a server-owned directory enumeration."""
        if not supplied or "\x00" in supplied:
            raise ValueError("invalid handoff note path")
        normalized = supplied.replace("\\", "/")
        windows_path = PureWindowsPath(supplied)
        supplied_is_absolute = bool(
            normalized.startswith("/") or windows_path.drive or windows_path.root
        )
        if supplied_is_absolute:
            requested_name = normalized.rsplit("/", 1)[-1]
            expected = directory.as_posix().rstrip("/") + "/" + requested_name
            if normalized != expected:
                raise ValueError("handoff note is outside its project directory")
        else:
            if "/" in normalized or normalized in {".", ".."}:
                raise ValueError("handoff note must be a filename")
            requested_name = normalized

        if HANDOFF_NOTE_RE.fullmatch(requested_name) is None:
            raise ValueError("handoff note is not a generated handoff document")
        resolved_directory = self._managed_path(directory)
        for entry in resolved_directory.iterdir():
            if entry.name != requested_name:
                continue
            if self._is_reparse_point(entry):
                raise ValueError("handoff note must not be a symbolic link")
            issued = self._managed_path(entry)
            if issued.parent != resolved_directory or not entry.is_file():
                raise ValueError("handoff note is not a regular project file")
            return issued
        raise FileNotFoundError(supplied)

    def _atomic_write(self, path: Path, content: str) -> None:
        destination = self._managed_path(path)
        parent = self._managed_path(destination.parent)
        parent.mkdir(parents=True, exist_ok=True)
        parent = self._managed_path(parent)
        destination = self._managed_path(path)
        if destination.parent != parent:
            raise ValueError("write destination escapes its managed directory")

        temporary = self.safe(parent, f".write-{uuid4().hex}.tmp")
        temporary = self._managed_path(temporary)
        try:
            with temporary.open("x", encoding="utf-8", newline="\n") as handle:
                handle.write(content)
                handle.flush()
                os.fsync(handle.fileno())
            # Revalidate every path immediately before the atomic rename. This
            # catches link/junction swaps that occurred while content was written.
            parent = self._managed_path(parent)
            destination = self._managed_path(path)
            temporary = self._managed_path(temporary)
            if destination.parent != parent or temporary.parent != parent:
                raise ValueError("atomic write paths escaped their managed directory")
            temporary.replace(destination)
        finally:
            with suppress(OSError, ValueError):
                cleanup = self._managed_path(temporary)
                cleanup.unlink()

    def _write_note(self, directory: Path, label: str, content: str) -> Path:
        if label not in NOTE_LABELS:
            raise ValueError("invalid note label")
        directory = self._managed_path(directory)
        directory.mkdir(parents=True, exist_ok=True)
        directory = self._managed_path(directory)
        stamp = datetime.now().strftime("%Y-%m-%d_%H%M%S_%f")
        note = self._managed_path(
            directory / f"{stamp}_{label}_{uuid4().hex[:6]}.md"
        )
        self._atomic_write(note, content)
        return note

    def _display_path(self, path: Path) -> str:
        return path.as_posix()

    @staticmethod
    def _collection_name(project_id: str) -> str:
        return "project_" + project_id.replace("-", "_")

    def _get_collection(self, name: str) -> Any | None:
        try:
            return self.client.get_collection(name)
        except Exception as exc:
            if "does not exist" in str(exc).lower() or "not found" in str(exc).lower():
                return None
            raise

    def _recover_transient_collections(self) -> None:
        """Recover an interrupted name swap and remove logical temporary collections."""
        with self.state_lock:
            try:
                projects = self._read_projects_unlocked()
                collections = list(self.client.list_collections())
            except Exception:
                return
            by_project: dict[str, list[Any]] = {}
            for collection in collections:
                metadata = collection.metadata or {}
                project_id = str(metadata.get("project_id") or "")
                if project_id:
                    by_project.setdefault(project_id, []).append(collection)
            registered = {project["id"]: project for project in projects}
            for project_id, project in registered.items():
                stable_name = self._collection_name(project_id)
                candidates = by_project.get(project_id, [])
                stable = next(
                    (item for item in candidates if item.name == stable_name), None
                )
                backups = [
                    item for item in candidates if item.name.startswith("backup_")
                ]
                expected = (project.get("index_manifest") or {}).get("fingerprint")
                matching_backup = next(
                    (
                        item
                        for item in backups
                        if not expected
                        or (item.metadata or {}).get("fingerprint") == expected
                    ),
                    None,
                )
                try:
                    if stable is None and matching_backup is not None:
                        matching_backup.modify(name=stable_name)
                        stable = matching_backup
                    elif (
                        stable is not None
                        and expected
                        and (stable.metadata or {}).get("fingerprint") != expected
                        and matching_backup is not None
                    ):
                        rejected = f"build_rejected_{project_id[:20]}_{uuid4().hex[:8]}"
                        stable.modify(name=rejected)
                        matching_backup.modify(name=stable_name)
                        self.client.delete_collection(rejected)
                except Exception:
                    continue
            with suppress(Exception):
                for collection in list(self.client.list_collections()):
                    if collection.name.startswith(
                        ("build_", "backup_", "build_rejected_")
                    ):
                        project_id = str(
                            (collection.metadata or {}).get("project_id") or ""
                        )
                        stable_name = (
                            self._collection_name(project_id) if project_id else ""
                        )
                        if stable_name and self._get_collection(stable_name) is None:
                            continue
                        with suppress(Exception):
                            self.client.delete_collection(collection.name)

    def cleanup_orphan_segments(self) -> dict[str, Any]:
        """Remove only old UUID segment directories absent from Chroma's catalog.

        Chroma 1.5 can leave HNSW directories behind after delete_collection.  The
        SQLite ``segments`` table is the authoritative set.  Any schema, locking,
        path or deletion error is reported as a warning and never affects serving.
        """
        state: dict[str, Any] = {
            "last_run_at": iso_now(),
            "removed_directories": 0,
            "removed_bytes": 0,
            "skipped_recent": 0,
            "warnings": [],
        }
        chroma_root = (self.data / "chroma").resolve()
        database = chroma_root / "chroma.sqlite3"
        if not database.is_file():
            self.cleanup_state = state
            return dict(state)

        try:
            uri = database.as_uri() + "?mode=ro"
            with sqlite3.connect(uri, uri=True, timeout=5) as connection:
                referenced = {
                    str(row[0]).lower()
                    for row in connection.execute("SELECT id FROM segments")
                }
        except Exception as exc:
            state["warnings"].append(f"unable to read Chroma segments: {exc}")
            self.cleanup_state = state
            return dict(state)

        cutoff = utc_now().timestamp() - self.orphan_grace_seconds
        try:
            candidates = list(chroma_root.iterdir())
        except OSError as exc:
            state["warnings"].append(f"unable to enumerate Chroma root: {exc}")
            self.cleanup_state = state
            return dict(state)

        with self.state_lock:
            for candidate in candidates:
                try:
                    # Never follow links or operate outside the exact Chroma root.
                    if candidate.is_symlink() or not candidate.is_dir():
                        continue
                    resolved = candidate.resolve()
                    if resolved.parent != chroma_root:
                        continue
                    try:
                        canonical = str(UUID(candidate.name)).lower()
                    except ValueError:
                        continue
                    if canonical != candidate.name.lower() or canonical in referenced:
                        continue
                    if candidate.stat().st_mtime > cutoff:
                        state["skipped_recent"] += 1
                        continue
                    size = sum(
                        path.stat().st_size
                        for path in candidate.rglob("*")
                        if path.is_file() and not path.is_symlink()
                    )
                    # Re-resolve immediately before deletion to close a path-swap
                    # window and keep the recursive target strictly scoped.
                    if (
                        candidate.is_symlink()
                        or candidate.resolve().parent != chroma_root
                    ):
                        continue
                    shutil.rmtree(candidate)
                    state["removed_directories"] += 1
                    state["removed_bytes"] += size
                except Exception as exc:
                    state["warnings"].append(f"{candidate.name}: {str(exc)[:1000]}")
        self.cleanup_state = state
        return dict(state)

    def watch_paths(self) -> list[Path]:
        managed_root = self.managed_root()
        managed_root.mkdir(parents=True, exist_ok=True)
        paths: list[Path] = [managed_root]
        for project in self.projects():
            for _, root in self._project_roots(project):
                if managed_root == root or managed_root in root.parents:
                    continue
                if root.exists() and root not in paths:
                    paths.append(root)
        return sorted(paths, key=str)

    def watch_signature(self) -> tuple[str, ...]:
        return tuple(str(path) for path in self.watch_paths())

    def affected_projects(self, changes: set[tuple[Any, str]]) -> set[str]:
        changed = [Path(path).resolve() for _, path in changes]
        affected: set[str] = set()
        for project in self.projects():
            for _, root in self._project_roots(project):
                if any(root == path or root in path.parents for path in changed):
                    affected.add(project["id"])
                    break
        return affected

    def _fast_source_reasons(self, project: dict[str, Any]) -> list[str]:
        """Detect source drift without re-hashing every file on each health poll.

        The five-minute consistency sweep remains the authoritative SHA-256 check.
        Health requests compare the persisted manifest with the current supported-file
        inventory (path, size and nanosecond mtime), which makes offline and live drift
        visible immediately without turning every connected UI's ten-second poll into a
        full-content scan.
        """
        manifest = project.get("index_manifest") or {}
        manifest_files = manifest.get("files")
        if not isinstance(manifest_files, list):
            return []
        expected = {
            str(item.get("source")): item
            for item in manifest_files
            if isinstance(item, dict) and item.get("source")
        }
        current: set[str] = set()
        added = changed = unreadable = False
        for kind, root in self._project_roots(project):
            if not root.exists():
                # A never-created managed overlay is intentionally absent.
                if kind == "overlay":
                    continue
                continue
            for path in self._iter_files(root):
                if path.suffix.lower() not in SUPPORTED_EXTENSIONS:
                    continue
                source = self._display_path(path)
                current.add(source)
                previous = expected.get(source)
                if previous is None:
                    added = True
                    continue
                try:
                    stat = path.stat()
                except OSError:
                    unreadable = True
                    continue
                if int(previous.get("size", -1)) != stat.st_size or (
                    previous.get("mtime_ns") is not None
                    and int(previous["mtime_ns"]) != stat.st_mtime_ns
                ):
                    changed = True
        reasons: list[str] = []
        if added:
            reasons.append("source files added")
        if set(expected) - current:
            reasons.append("source files removed")
        if changed:
            reasons.append("source files changed")
        if unreadable:
            reasons.append("source inventory unreadable")
        return reasons

    def health_report(self) -> dict[str, Any]:
        projects = self.projects()
        source_reasons = {
            project["id"]: self._fast_source_reasons(project) for project in projects
        }
        indexed = 0
        stale = 0
        total_chunks = 0
        failed_files = 0
        truncated_files = 0
        unsupported_files = 0
        missing_roots = 0
        items: list[dict[str, Any]] = []
        with self.state_lock:
            for project in projects:
                expected = int(project.get("chunks") or 0)
                total_chunks += expected
                collection_error: str | None = None
                try:
                    collection = self._get_collection(
                        self._collection_name(project["id"])
                    )
                    count = collection.count() if collection is not None else None
                except Exception as exc:
                    collection = None
                    count = None
                    collection_error = str(exc)[:1000]
                manifest = project.get("index_manifest") or {}
                stats = project.get("index_stats") or {}
                roots_missing = sum(
                    not root.exists()
                    for kind, root in self._project_roots(project)
                    if kind != "overlay"
                )
                missing_roots += roots_missing
                failed_files += int(stats.get("failed_files") or 0)
                truncated_files += len(stats.get("truncated_files") or [])
                unsupported_files += int(stats.get("unsupported_files") or 0)
                reasons: list[str] = []
                if collection_error:
                    reasons.append("collection unreadable")
                elif collection is None:
                    reasons.append("collection missing")
                elif count != expected:
                    reasons.append("collection count mismatch")
                else:
                    indexed += 1
                if not manifest:
                    reasons.append("manifest missing")
                if manifest.get("embedding_model") != self.embedding_model:
                    reasons.append("embedding model changed")
                if manifest.get("embedding_revision") != self.embedding_revision:
                    reasons.append("embedding revision changed")
                if manifest.get("chunk_version") != CHUNK_VERSION:
                    reasons.append("chunk version changed")
                if manifest.get("max_text_chars") != self.max_text_chars:
                    reasons.append("text limit changed")
                if stats.get("last_error"):
                    reasons.append("last index attempt failed")
                if roots_missing:
                    reasons.append("source root missing")
                reasons.extend(source_reasons.get(project["id"], []))
                reasons = list(dict.fromkeys(reasons))
                if reasons:
                    stale += 1
                items.append(
                    {
                        "id": project["id"],
                        "chunks": expected,
                        "collection_chunks": count,
                        "indexed_at": project.get("indexed_at"),
                        "fingerprint": manifest.get("fingerprint"),
                        "stale": bool(reasons),
                        "reasons": reasons,
                        "failed_files": int(stats.get("failed_files") or 0),
                        "truncated_files": len(stats.get("truncated_files") or []),
                        "unsupported_files": int(stats.get("unsupported_files") or 0),
                    }
                )
        with self.state_lock:
            watcher_failed = bool(self.watcher_state.get("last_error"))
            watcher = {
                "running": bool(self.watcher_state.get("running")),
                "roots": int(self.watcher_state.get("roots") or 0),
                "restarts": int(self.watcher_state.get("restarts") or 0),
                "indexing_projects": sorted(
                    self.watcher_state.get("indexing_projects") or []
                ),
                "last_event_at": self.watcher_state.get("last_event_at"),
                "last_error": "watcher operation failed" if watcher_failed else None,
                "last_recovered_at": self.watcher_state.get("last_recovered_at"),
            }
            sweep_failed = bool(self.sweep_state.get("last_error"))
            sweep = {
                "running": bool(self.sweep_state.get("running")),
                "last_started_at": self.sweep_state.get("last_started_at"),
                "last_completed_at": self.sweep_state.get("last_completed_at"),
                "last_error": "consistency sweep failed" if sweep_failed else None,
            }
            background_failed_raw = bool(
                self.background_index_state.get("last_error")
            )
            background = {
                "last_started_at": self.background_index_state.get(
                    "last_started_at"
                ),
                "last_completed_at": self.background_index_state.get(
                    "last_completed_at"
                ),
                "last_error": (
                    "background index failed" if background_failed_raw else None
                ),
                "last_error_at": self.background_index_state.get("last_error_at"),
                "running_projects": sorted(self._background_index_running),
                "queued_projects": sorted(
                    project_id
                    for project_id, reasons in self._background_index_pending.items()
                    if reasons
                ),
            }
            cleanup_warnings = len(self.cleanup_state.get("warnings") or [])
            cleanup = {
                "last_run_at": self.cleanup_state.get("last_run_at"),
                "removed_directories": int(
                    self.cleanup_state.get("removed_directories") or 0
                ),
                "removed_bytes": int(self.cleanup_state.get("removed_bytes") or 0),
                "skipped_recent": int(self.cleanup_state.get("skipped_recent") or 0),
                "warnings": (
                    ["cleanup operation reported a warning"] * cleanup_warnings
                ),
            }
        watcher_busy = bool(watcher.get("indexing_projects"))
        sweep_ready = bool(sweep.get("last_completed_at"))
        sweep_busy = bool(sweep.get("running"))
        background_busy = bool(
            background.get("running_projects") or background.get("queued_projects")
        )
        background_failed = background_failed_raw and any(
            "last index attempt failed" in item["reasons"] for item in items
        )
        healthy = (
            stale == 0
            and watcher.get("running", False)
            and not watcher.get("last_error")
            and not watcher_busy
            and sweep_ready
            and not sweep_busy
            and not sweep_failed
            and not background_busy
            and not background_failed
        )
        return {
            "status": "ok" if healthy else "degraded",
            "watcher": watcher,
            "sweep": sweep,
            "background_index": background,
            "cleanup": cleanup,
            "index": {
                "embedding_model": self.embedding_model,
                "embedding_revision": self.embedding_revision,
                "chunk_version": CHUNK_VERSION,
                "registered_projects": len(projects),
                "indexed_projects": indexed,
                "stale_projects": stale,
                "chunks": total_chunks,
                "failed_files": failed_files,
                "truncated_files": truncated_files,
                "unsupported_files": unsupported_files,
                "missing_roots": missing_roots,
                "projects": items,
            },
        }


service: KnowledgeService | None = None


async def watch_vault() -> None:
    assert service is not None
    while True:
        try:
            paths = service.watch_paths()
            signature = tuple(str(path) for path in paths)
            service.watcher_state.update(
                {"running": True, "roots": len(paths), "last_error": None}
            )
            async for changes in awatch(
                *paths,
                debounce=1500,
                rust_timeout=WATCH_REFRESH_SECONDS * 1000,
                yield_on_timeout=True,
                force_polling=True,
                poll_delay_ms=WATCH_POLL_MS,
                ignore_permission_denied=True,
            ):
                if tuple(str(path) for path in service.watch_paths()) != signature:
                    break
                if not changes:
                    continue
                service.watcher_state["last_event_at"] = iso_now()
                for project_id in service.affected_projects(changes):
                    with service.state_lock:
                        indexing = set(
                            service.watcher_state.get("indexing_projects") or []
                        )
                        indexing.add(project_id)
                        service.watcher_state["indexing_projects"] = sorted(indexing)
                    try:
                        await service.index_if_stale(
                            project_id, reason="filesystem event"
                        )
                    except Exception as exc:
                        service._record_index_error(project_id, exc)
                    finally:
                        with service.state_lock:
                            indexing = set(
                                service.watcher_state.get("indexing_projects") or []
                            )
                            indexing.discard(project_id)
                            service.watcher_state["indexing_projects"] = sorted(
                                indexing
                            )
            service.watcher_state["last_recovered_at"] = iso_now()
        except asyncio.CancelledError:
            service.watcher_state["running"] = False
            raise
        except Exception as exc:
            service.watcher_state.update(
                {
                    "running": False,
                    "last_error": str(exc)[:4000],
                    "restarts": int(service.watcher_state.get("restarts") or 0) + 1,
                }
            )
            await asyncio.sleep(2)


async def consistency_sweep() -> None:
    assert service is not None
    first = True
    while True:
        try:
            if not first:
                await asyncio.sleep(SWEEP_SECONDS)
            first = False
            service.sweep_state.update(
                {
                    "running": True,
                    "last_started_at": iso_now(),
                    "last_error": None,
                }
            )
            results = await service.index_stale(
                reason=(
                    "startup"
                    if service.sweep_state["last_completed_at"] is None
                    else "periodic sweep"
                )
            )
            service.cleanup_orphan_segments()
            errors = [item["error"] for item in results if item.get("error")]
            service.sweep_state.update(
                {
                    "running": False,
                    "last_completed_at": iso_now(),
                    "last_error": "; ".join(errors[:3]) if errors else None,
                }
            )
        except asyncio.CancelledError:
            service.sweep_state["running"] = False
            raise
        except Exception as exc:
            service.sweep_state.update(
                {"running": False, "last_error": str(exc)[:4000]}
            )
            await asyncio.sleep(5)


@asynccontextmanager
async def lifespan(_: FastAPI):
    global service
    service = KnowledgeService()
    watcher = asyncio.create_task(watch_vault(), name="knowledge-watcher")
    sweeper = asyncio.create_task(consistency_sweep(), name="knowledge-sweeper")
    await asyncio.sleep(0)
    try:
        yield
    finally:
        watcher.cancel()
        sweeper.cancel()
        await asyncio.gather(watcher, sweeper, return_exceptions=True)
        await service.cancel_background_indexes()


app = FastAPI(title="Local-First Agent Knowledge Service", lifespan=lifespan)


@app.get("/health")
def health() -> dict[str, Any]:
    assert service is not None
    return service.health_report()


@app.get("/projects")
def projects(
    include_archived: bool = FastAPIQuery(default=False),
) -> list[dict[str, Any]]:
    assert service is not None
    items = service.projects()
    if not include_archived:
        items = [item for item in items if not item.get("archived")]
    return [
        {
            **service.public_project(item),
            "archived": bool(item.get("archived")),
        }
        for item in items
    ]


@app.get("/projects/{project_id}")
def project(project_id: str) -> dict[str, Any]:
    assert service is not None
    try:
        return service.describe(project_id)
    except KeyError as exc:
        raise HTTPException(404, "project was not found") from exc


@app.post("/projects")
async def create(project: ProjectCreate) -> dict[str, Any]:
    assert service is not None
    try:
        service.add(project)
        await service.index(project.id, reason="project creation")
        return service.describe(project.id)
    except ValueError as exc:
        raise HTTPException(409, str(exc)) from exc


@app.patch("/projects/{project_id}")
def update_project(project_id: str, request: ProjectUpdate) -> dict[str, Any]:
    assert service is not None
    if request.name is None and request.archived is None:
        raise HTTPException(422, "at least one field must be provided")
    try:
        return service.update_project(project_id, request)
    except KeyError as exc:
        raise HTTPException(404, "project was not found") from exc


@app.delete("/projects/{project_id}")
async def delete_project(project_id: str, request: ProjectDelete) -> dict[str, Any]:
    assert service is not None
    try:
        return await service.delete_project(project_id, request)
    except KeyError as exc:
        raise HTTPException(404, "project was not found") from exc
    except ValueError as exc:
        raise HTTPException(422, str(exc)) from exc


@app.post("/projects/{project_id}/search")
def search(project_id: str, request: Query) -> dict[str, Any]:
    assert service is not None
    try:
        return service.search(project_id, request)
    except KeyError as exc:
        raise HTTPException(404, "project has not been indexed") from exc


@app.post("/search")
def search_all(request: Query) -> dict[str, Any]:
    assert service is not None
    return service.search_all(request)


@app.post("/projects/{project_id}/index")
async def index(project_id: str) -> dict[str, Any]:
    assert service is not None
    try:
        return await service.index(project_id, reason="manual rebuild")
    except KeyError as exc:
        raise HTTPException(404, "project has not been indexed") from exc


@app.post("/projects/{project_id}/research")
async def research(project_id: str, request: ResearchCapture) -> dict[str, Any]:
    assert service is not None
    try:
        return await service.capture_research(project_id, request)
    except KeyError as exc:
        raise HTTPException(404, "project has not been indexed") from exc


@app.post("/projects/{project_id}/progress")
async def progress(project_id: str, request: ProgressCapture) -> dict[str, Any]:
    assert service is not None
    try:
        return await service.capture_progress(project_id, request)
    except KeyError as exc:
        raise HTTPException(404, "project has not been indexed") from exc


@app.post("/projects/{project_id}/handoff")
async def handoff(project_id: str, request: CodexHandoff) -> dict[str, Any]:
    assert service is not None
    try:
        return await service.capture_handoff(project_id, request)
    except KeyError as exc:
        raise HTTPException(404, "project has not been indexed") from exc


@app.post("/projects/{project_id}/handoff-result")
async def handoff_result(
    project_id: str, request: CodexHandoffResult
) -> dict[str, Any]:
    assert service is not None
    try:
        return await service.capture_handoff_result(project_id, request)
    except KeyError as exc:
        raise HTTPException(404, "project has not been indexed") from exc
    except FileNotFoundError as exc:
        raise HTTPException(404, "handoff note was not found") from exc
